"""Execute :class:`VerificationStep` specs against a workspace directory."""

import ast
import json
import logging
import os
import re
import shlex
import signal
import socket
import subprocess
import sys
import threading
import time
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import replace
from pathlib import Path
from typing import IO, Any

from collections.abc import Iterable

try:  # Python 3.11+
    import tomllib as _toml  # type: ignore[import-not-found]
except ModuleNotFoundError:  # pragma: no cover - fallback for <3.11
    import tomli as _toml  # type: ignore[no-redef]

from umbrella.verification.models import (
    VerificationReport,
    VerificationStatus,
    VerificationStep,
    VerificationStepKind,
    VerificationStepResult,
)
from umbrella.deep_agent_tools.domain_policy import public_workspace_llm_env_bridge
from umbrella.enforcement import (
    check_verification_step_diff,
    diff_snapshots,
    snapshot_workspace,
)
from umbrella.verification.skill_compliance import build_skill_compliance_results
from umbrella.verification.source_policy import (
    mock_scaffold_hits,
    scan_changed_files_for_mock_scaffold,
)
from umbrella.verification.spec_loader import load_verification_meta
from umbrella.verification.diff_policy import run_diff_policy_guard
from umbrella.verification.mutation_smoke import run_mutation_smoke_guard
from umbrella.verification.test_quality import run_test_quality_guard

log = logging.getLogger(__name__)

_DEFAULT_OVERALL_TIMEOUT = 900  # 15 minutes hard ceiling

_EXECUTABLE_STEP_KINDS: frozenset[VerificationStepKind] = frozenset(
    {
        VerificationStepKind.SHELL,
        VerificationStepKind.FILE_EXISTS,
        VerificationStepKind.HTTP_BOOT,
        VerificationStepKind.IMPORT_CHECK,
        VerificationStepKind.BEHAVIORAL_HTTP,
        VerificationStepKind.INPUT_SENSITIVITY,
        VerificationStepKind.PPTX_DIFF,
        VerificationStepKind.SOURCE_POLICY,
    }
)


def _load_detected_domains_for_workspace(workspace_id: str) -> set[str]:
    """Best-effort lookup of the skill layer's verdict for ``workspace_id``.

    Reads ``.umbrella/ouroboros_drive/state/active_skills.json`` relative to
    the current working directory of the runner host process. Returns an
    empty set on any error -- compliance checks are then skipped entirely
    rather than failing the run on a missing/stale cache.
    """
    try:
        repo_root = Path.cwd()
        domains_path = (
            repo_root / "workspaces" / workspace_id / ".memory" / "domains.json"
        )
        payload = json.loads(domains_path.read_text(encoding="utf-8"))
        raw = payload.get("domains") if isinstance(payload, dict) else None
        if isinstance(raw, list):
            domains = {str(value) for value in raw if str(value).strip()}
            return _apply_domain_overrides_for_workspace_id(domains, workspace_id)
    except Exception:
        pass
    try:
        from umbrella.orchestration.ouroboros_task import load_detected_domains

        domains = load_detected_domains(Path.cwd(), workspace_id)
        return _apply_domain_overrides_for_workspace_id(domains, workspace_id)
    except Exception:  # noqa: BLE001 - best-effort, must never break verify
        log.debug("Failed to load detected domains", exc_info=True)
        return set()


def _apply_domain_overrides_for_workspace_id(
    detected_domains: set[str], workspace_id: str
) -> set[str]:
    workspace_path = Path.cwd() / "workspaces" / workspace_id
    return _apply_domain_overrides(detected_domains, workspace_path)


def _apply_domain_overrides(
    detected_domains: set[str], workspace_path: Path
) -> set[str]:
    """Apply explicit workspace.toml policy on top of detected domains."""

    domains = {str(value) for value in detected_domains if str(value).strip()}
    if _workspace_declares_gmas_disabled(workspace_path):
        domains.discard("multi_agent_gmas")
    return domains


def _workspace_declares_gmas_disabled(workspace_path: Path) -> bool:
    """Return ``True`` if ``workspace.toml`` opts out of ``multi_agent_gmas``.

    Note: this is a defensive helper — the verification runner no longer
    emits any synthetic ``skill_compliance:multi_agent_gmas*`` checks
    (``build_skill_compliance_results`` is a no-op now), so this function
    is kept only for backward compatibility with callers that read the
    explicit policy.
    """
    path = workspace_path / "workspace.toml"
    try:
        with path.open("rb") as fh:
            data: dict[str, Any] = _toml.load(fh)
    except Exception:
        return False

    skills = data.get("skills")
    if isinstance(skills, dict) and skills.get("multi_agent_gmas") is False:
        return True

    gmas = data.get("gmas")
    if isinstance(gmas, dict) and gmas.get("enabled") is False:
        return True

    workspace = data.get("workspace")
    if isinstance(workspace, dict):
        if workspace.get("requires_gmas") is False:
            return True
        if workspace.get("multi_agent_gmas") is False:
            return True

    return False


def run_verification(
    workspace_path: str | Path,
    steps: Iterable[VerificationStep],
    *,
    workspace_id: str | None = None,
    overall_timeout_seconds: int = _DEFAULT_OVERALL_TIMEOUT,
    env: dict[str, str] | None = None,
    detected_domains: set[str] | None = None,
    changed_files: list[str] | None = None,
) -> VerificationReport:
    """Run ``steps`` in order against ``workspace_path`` and return a report.

    The runner is intentionally defensive: a crashing step is converted to
    ``VerificationStatus.ERROR`` rather than propagating.  We stop running
    subsequent steps only when an overall timeout is breached, otherwise
    every step is attempted (so the retry prompt has a full picture).
    """

    workspace_path = Path(workspace_path).resolve()
    workspace_id = workspace_id or workspace_path.name
    started = time.time()
    results: list[VerificationStepResult] = []
    changed_files_param = list(changed_files) if changed_files else []

    base_env = os.environ.copy()
    if env:
        base_env.update(env)
    base_env.update(public_workspace_llm_env_bridge(base_env))
    # Unbuffered output keeps child logs legible.
    base_env.setdefault("PYTHONUNBUFFERED", "1")
    base_env.setdefault("PYTHONIOENCODING", "utf-8")

    vmeta = load_verification_meta(workspace_path)
    skip_tq = bool(vmeta.get("skip_test_quality", False))
    chosen_python = _choose_python_executable(workspace_path)
    log.info(
        "Verification interpreter for %s: %s", workspace_path, " ".join(chosen_python)
    )

    filtered_steps = [step for step in steps if step.kind in _EXECUTABLE_STEP_KINDS]
    code_task = _changed_files_indicate_code_task(changed_files_param)
    ran_test_quality = False

    for step in filtered_steps:
        remaining = overall_timeout_seconds - (time.time() - started)
        if remaining <= 0:
            results.append(
                VerificationStepResult(
                    step=step,
                    status=VerificationStatus.SKIPPED,
                    summary="Skipped: overall verification timeout exhausted",
                )
            )
            continue

        step_env = dict(base_env)
        if step.env:
            step_env.update(step.env)
            step_env.update(public_workspace_llm_env_bridge(step_env))

        effective_step = _with_workspace_command(step, chosen_python)
        before_step_snapshot = snapshot_workspace(workspace_path)

        try:
            if effective_step.kind == VerificationStepKind.SHELL:
                result = _run_shell_step(effective_step, workspace_path, step_env)
            elif effective_step.kind == VerificationStepKind.FILE_EXISTS:
                result = _run_file_exists_step(effective_step, workspace_path)
            elif effective_step.kind == VerificationStepKind.IMPORT_CHECK:
                result = _run_import_step(effective_step, workspace_path, step_env)
            elif effective_step.kind == VerificationStepKind.HTTP_BOOT:
                result = _run_http_boot_step(effective_step, workspace_path, step_env)
            elif effective_step.kind == VerificationStepKind.BEHAVIORAL_HTTP:
                result = _run_behavioral_http_step(
                    effective_step, workspace_path, step_env
                )
            elif effective_step.kind == VerificationStepKind.INPUT_SENSITIVITY:
                result = _run_input_sensitivity_step(
                    effective_step, workspace_path, step_env
                )
            elif effective_step.kind == VerificationStepKind.PPTX_DIFF:
                result = _run_pptx_diff_step(effective_step, workspace_path, step_env)
            elif effective_step.kind == VerificationStepKind.SOURCE_POLICY:
                scan_files = (
                    changed_files_param
                    if changed_files_param
                    else _workspace_policy_scan_files(workspace_path)
                )
                result = _run_source_policy_step(
                    workspace_path=workspace_path,
                    workspace_id=workspace_id,
                    changed_files=scan_files,
                    step=effective_step,
                )
            else:
                result = VerificationStepResult(
                    step=effective_step,
                    status=VerificationStatus.ERROR,
                    error=f"Unknown verification step kind: {effective_step.kind!r}",
                )
        except Exception as exc:  # noqa: BLE001 - defensive boundary
            log.exception("Verification step %s crashed", step.name)
            result = VerificationStepResult(
                step=step,
                status=VerificationStatus.ERROR,
                error=f"{type(exc).__name__}: {exc}",
            )

        results.append(result)
        verifier_changes = diff_snapshots(
            before_step_snapshot,
            snapshot_workspace(workspace_path),
        )
        verifier_issues = check_verification_step_diff(verifier_changes)
        if verifier_issues:
            results.append(
                VerificationStepResult(
                    step=VerificationStep(
                        kind=VerificationStepKind.SOURCE_POLICY,
                        name="enforcement:verification_mutation",
                        optional=False,
                    ),
                    status=VerificationStatus.FAILED,
                    summary="\n".join(
                        f"{issue.code}: {issue.path} {issue.message}"
                        for issue in verifier_issues[:20]
                    ),
                    error="verification_step_mutated_candidate_evaluator_boundary",
                )
            )

        if (
            not skip_tq
            and step.kind == VerificationStepKind.SHELL
            and _command_looks_like_pytest(effective_step.command)
        ):
            gres = run_test_quality_guard(workspace_path, require_tests=code_task)
            results.append(gres)
            ran_test_quality = True

    ran_source_in_loop = any(
        s.kind == VerificationStepKind.SOURCE_POLICY for s in filtered_steps
    )
    if changed_files_param and not ran_source_in_loop:
        remaining = overall_timeout_seconds - (time.time() - started)
        if remaining <= 0:
            results.append(
                VerificationStepResult(
                    step=VerificationStep(
                        kind=VerificationStepKind.SOURCE_POLICY,
                        name="source_policy:changed_files_auto",
                        optional=False,
                    ),
                    status=VerificationStatus.SKIPPED,
                    summary="Skipped: overall verification timeout exhausted",
                )
            )
        else:
            results.append(
                _run_source_policy_step(
                    workspace_path=workspace_path,
                    workspace_id=workspace_id,
                    changed_files=changed_files_param,
                    step=None,
                )
            )
    if code_task and not skip_tq and not ran_test_quality:
        results.append(run_test_quality_guard(workspace_path, require_tests=True))

    if code_task:
        results.append(
            run_diff_policy_guard(
                workspace_path,
                changed_files=changed_files_param,
                approved_policy_edits=False,
            )
        )
        results.append(
            run_mutation_smoke_guard(
                workspace_path,
                changed_files=changed_files_param,
                python_cmd=chosen_python,
                env=base_env,
            )
        )

    if detected_domains is None:
        detected_domains = _load_detected_domains_for_workspace(workspace_id)
    else:
        detected_domains = _apply_domain_overrides(detected_domains, workspace_path)
    if detected_domains:
        for compliance_result in build_skill_compliance_results(
            workspace_path,
            detected_domains,
            python_cmd=chosen_python,
            env=base_env,
        ):
            results.append(compliance_result)

    results.append(_run_python_compile_check(workspace_path))

    return VerificationReport(
        workspace_id=workspace_id,
        workspace_path=str(workspace_path),
        results=results,
        started_at=started,
        finished_at=time.time(),
    )


_CODE_TASK_EXTS = {".py", ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs", ".go", ".rs"}


def _changed_files_indicate_code_task(changed_files: Iterable[str]) -> bool:
    for raw in changed_files:
        path = str(raw or "").replace("\\", "/").strip().lstrip("./")
        if not path:
            continue
        parts = [part.lower() for part in path.split("/") if part]
        if not parts:
            continue
        if parts[0] in {"tests", "test", "docs", "doc", ".memory", ".umbrella"}:
            continue
        name = parts[-1]
        if name.startswith("test_") and name.endswith(".py"):
            continue
        if Path(path).suffix.lower() in _CODE_TASK_EXTS:
            return True
    return False


def _run_python_compile_check(workspace_path: Path) -> VerificationStepResult:
    """Fail verification when Python source has syntax/import-time compile errors."""

    step = VerificationStep(
        kind=VerificationStepKind.IMPORT_CHECK,
        name="python_compile_check",
        optional=False,
    )
    start = time.time()
    cmd = [sys.executable, "-m", "compileall", "-q", str(workspace_path)]
    try:
        proc = subprocess.run(
            cmd,
            cwd=str(workspace_path),
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=120,
            check=False,
        )
    except Exception as exc:  # noqa: BLE001
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            duration_seconds=time.time() - start,
            summary=f"compile check crashed: {type(exc).__name__}",
            error=f"{type(exc).__name__}: {exc}",
        )
    status = (
        VerificationStatus.PASSED if proc.returncode == 0 else VerificationStatus.FAILED
    )
    return VerificationStepResult(
        step=step,
        status=status,
        exit_code=proc.returncode,
        duration_seconds=time.time() - start,
        stdout=proc.stdout or "",
        stderr=proc.stderr or "",
        summary=f"python -m compileall -q {workspace_path} -> exit {proc.returncode}",
        error="" if proc.returncode == 0 else "python_compile_failed",
    )


def _choose_python_executable(workspace_path: Path) -> list[str]:
    venv_py = (
        workspace_path
        / ".venv"
        / ("Scripts/python.exe" if sys.platform == "win32" else "bin/python")
    )
    if venv_py.exists():
        return [str(venv_py)]
    repo_root = None
    try:
        parts = workspace_path.resolve().parts
        if "workspaces" in parts:
            idx = len(parts) - 1 - list(reversed(parts)).index("workspaces")
            if idx > 0:
                repo_root = Path(*parts[:idx])
    except Exception:
        repo_root = None
    if repo_root is not None:
        repo_venv_py = (
            repo_root
            / ".venv"
            / ("Scripts/python.exe" if sys.platform == "win32" else "bin/python")
        )
        if repo_venv_py.exists():
            return [str(repo_venv_py)]
    if (workspace_path / "pyproject.toml").exists():
        return ["uv", "run", "python"]
    return [sys.executable]


_PYTHON_NAMES = {"python", "python3", "py"}
_SHELL_NAMES = {"bash", "sh"}
_SHELL_C_FLAGS = {"-c", "-lc"}
_SHELL_META_TOKENS = {"&&", "||", ";", "|", ">", ">>", "<", "2>", "2>>", "&"}


def _with_workspace_command(
    step: VerificationStep,
    python_cmd: list[str],
) -> VerificationStep:
    """Return ``step`` with command argv normalized for this workspace.

    Explicit workspace specs are often authored by an LLM. Keep the public
    spec permissive, but make execution deterministic: use the selected
    workspace interpreter, unwrap simple ``bash -c 'python ...'`` wrappers,
    and repair a common over-escaped ``python -c`` payload only when the
    original source is invalid and the repaired source parses.
    """

    command = _normalize_verification_command(step.command, python_cmd)
    if step.kind == VerificationStepKind.IMPORT_CHECK and not command and step.module:
        command = [*python_cmd, "-c", f"import {step.module}"]
    if command == step.command:
        return step
    return replace(step, command=command)


def _normalize_verification_command(
    command: Iterable[str],
    python_cmd: list[str],
) -> list[str]:
    cmd = [str(part) for part in command]
    if not cmd:
        return []
    unwrapped = _unwrap_simple_shell_python_command(cmd)
    if unwrapped is not None:
        cmd = unwrapped
    cmd = _repair_python_c_args(cmd)
    return _rewrite_python_executable(cmd, python_cmd)


def _unwrap_simple_shell_python_command(cmd: list[str]) -> list[str] | None:
    if len(cmd) < 3:
        return None
    if cmd[0].lower() not in _SHELL_NAMES or cmd[1].lower() not in _SHELL_C_FLAGS:
        return None
    parsed = _split_shell_python_command(cmd[2])
    if not parsed or any(token in _SHELL_META_TOKENS for token in parsed):
        return None
    if _python_invocation_start(parsed) is None:
        return None
    return _repair_python_c_args(parsed)


def _split_shell_python_command(command: str) -> list[str]:
    try:
        parsed = shlex.split(command, posix=True)
    except ValueError:
        parsed = []
    if _python_invocation_start(parsed) is not None:
        return parsed
    if os.name == "nt":
        try:
            return shlex.split(command, posix=False)
        except ValueError:
            return []
    return parsed


def _python_invocation_start(cmd: list[str]) -> int | None:
    if not cmd:
        return None
    if _is_python_token(cmd[0]):
        return 0
    lowered = [part.lower() for part in cmd[:3]]
    if (
        len(cmd) >= 3
        and lowered[0] == "uv"
        and lowered[1] == "run"
        and _is_python_token(cmd[2])
    ):
        return 2
    return None


def _rewrite_python_executable(cmd: list[str], python_cmd: list[str]) -> list[str]:
    if not cmd:
        return cmd
    if _is_python_token(cmd[0]):
        return [*python_cmd, *cmd[1:]]
    lowered = [part.lower() for part in cmd[:3]]
    if (
        len(cmd) >= 3
        and lowered[0] == "uv"
        and lowered[1] == "run"
        and _is_python_token(cmd[2])
    ):
        return [*python_cmd, *cmd[3:]]
    return cmd


def _repair_python_c_args(cmd: list[str]) -> list[str]:
    source_idx = _python_c_source_index(cmd)
    if source_idx is None or source_idx >= len(cmd):
        return cmd
    source = cmd[source_idx]
    repaired = _repair_python_c_source(source)
    if repaired == source:
        return cmd
    updated = list(cmd)
    updated[source_idx] = repaired
    return updated


def _python_c_source_index(cmd: list[str]) -> int | None:
    if len(cmd) >= 3 and _is_python_token(cmd[0]) and cmd[1] == "-c":
        return 2
    lowered = [part.lower() for part in cmd[:4]]
    if (
        len(cmd) >= 5
        and lowered[0] == "uv"
        and lowered[1] == "run"
        and _is_python_token(cmd[2])
        and cmd[3] == "-c"
    ):
        return 4
    return None


def _is_python_token(token: str) -> bool:
    name = Path(str(token)).name.lower()
    return name in _PYTHON_NAMES or re.fullmatch(r"python3(?:\.\d+)?", name) is not None


def _repair_python_c_source(source: str) -> str:
    if '\\"' not in source and "\\'" not in source:
        return source
    try:
        ast.parse(source)
        return source
    except SyntaxError:
        pass
    candidates = [
        source.replace('\\"', '"'),
        source.replace("\\'", "'"),
        source.replace('\\"', '"').replace("\\'", "'"),
    ]
    for candidate in candidates:
        if candidate == source:
            continue
        try:
            ast.parse(candidate)
        except SyntaxError:
            continue
        return candidate
    return source


def _has_behavioral_step(results: list[VerificationStepResult]) -> bool:
    behavioral = {
        VerificationStepKind.BEHAVIORAL_HTTP,
        VerificationStepKind.INPUT_SENSITIVITY,
        VerificationStepKind.PPTX_DIFF,
        VerificationStepKind.SOURCE_POLICY,
    }
    return any(r.step.kind in behavioral for r in results)


def _run_source_policy_step(
    *,
    workspace_path: Path,
    workspace_id: str,
    changed_files: list[str],
    step: VerificationStep | None = None,
) -> VerificationStepResult:
    repo_root = workspace_path
    if (
        workspace_path.name == workspace_id
        and workspace_path.parent.name == "workspaces"
    ):
        repo_root = workspace_path.parent.parent
    hits = scan_changed_files_for_mock_scaffold(
        repo_root=repo_root,
        workspace_path=workspace_path,
        changed_files=changed_files,
    )
    return VerificationStepResult(
        step=step
        or VerificationStep(
            kind=VerificationStepKind.SOURCE_POLICY,
            name="source_policy:mock_scaffold_scan",
            optional=False,
        ),
        status=VerificationStatus.PASSED,
        summary=(
            "No mock/scaffold markers found in changed source files."
            if not hits
            else "Advisory mock/scaffold markers found in changed source files: "
            + "; ".join(hits[:10])
        ),
        error="",
        stdout="\n".join(hits),
    )


def _workspace_policy_scan_files(workspace_path: Path) -> list[str]:
    """Walk the workspace and return rel-paths the policy step should
    inspect. Combines hard-coded directory excludes with the
    user-configurable ``[verification.skip_paths]`` glob list so the
    full-scan and incremental modes stay consistent (otherwise a
    ``.memory/`` lesson file could be scanned by the full mode but
    skipped by the incremental ``changed_files`` mode, leading to
    flaky verification).
    """
    from umbrella.verification.source_policy import (
        _glob_matches_any,
        load_skip_path_patterns,
    )

    ignored_dirs = {
        ".git",
        ".venv",
        ".memory",
        ".umbrella",
        ".umbrella_scratch",
        "__pycache__",
        "node_modules",
    }
    skip_globs = load_skip_path_patterns(workspace_path)
    rels: list[str] = []
    for path in workspace_path.rglob("*"):
        if not path.is_file():
            continue
        rel = path.relative_to(workspace_path)
        if any(part in ignored_dirs for part in rel.parts):
            continue
        rel_posix = rel.as_posix()
        if _glob_matches_any(rel_posix, skip_globs):
            continue
        rels.append(rel_posix)
    return rels


def _run_file_exists_step(
    step: VerificationStep, workspace_path: Path
) -> VerificationStepResult:
    rel = (step.path or "").strip()
    if not rel:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="file_exists requires path",
        )
    target = (workspace_path / rel).resolve()
    try:
        target.relative_to(workspace_path.resolve())
    except ValueError:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error=f"file_exists path escapes workspace: {rel}",
        )
    if target.exists():
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.PASSED,
            summary=f"Found {rel}",
        )
    return VerificationStepResult(
        step=step,
        status=VerificationStatus.FAILED,
        summary=f"Missing required path: {rel}",
    )


# ---------------------------------------------------------------------------
# Shell step
# ---------------------------------------------------------------------------


def _run_shell_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    if not step.command:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="shell step requires a non-empty `command`",
        )

    start = time.time()
    try:
        proc = subprocess.run(  # noqa: S603 - intentional subprocess
            list(step.command),
            cwd=str(cwd),
            env=env,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=max(1, step.timeout_seconds),
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.FAILED,
            duration_seconds=time.time() - start,
            summary=f"Timed out after {step.timeout_seconds}s: {_fmt_cmd(step.command)}",
            stdout=exc.stdout or "" if isinstance(exc.stdout, str) else "",
            stderr=(exc.stderr or "") if isinstance(exc.stderr, str) else "",
            error="timeout",
        )
    except FileNotFoundError as exc:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            duration_seconds=time.time() - start,
            error=f"executable not found: {exc}",
            summary=f"Cannot run {_fmt_cmd(step.command)}: {exc}",
        )

    duration = time.time() - start
    stdout = proc.stdout or ""
    stderr = proc.stderr or ""
    status = (
        VerificationStatus.PASSED if proc.returncode == 0 else VerificationStatus.FAILED
    )
    combined_output = f"{stdout}\n{stderr}"
    if _command_looks_like_compileall(step.command) and "Can't list" in combined_output:
        status = VerificationStatus.FAILED
    return VerificationStepResult(
        step=step,
        status=status,
        exit_code=proc.returncode,
        duration_seconds=duration,
        stdout=stdout,
        stderr=stderr,
        summary=f"{_fmt_cmd(step.command)} -> exit {proc.returncode} in {duration:.1f}s",
        error=(
            "python_compile_failed_cannot_list"
            if _command_looks_like_compileall(step.command)
            and "Can't list" in combined_output
            else ""
        ),
    )


# ---------------------------------------------------------------------------
# HTTP boot step
# ---------------------------------------------------------------------------


def _command_looks_like_pytest(command: list[str]) -> bool:
    joined = " ".join(str(c) for c in command).lower()
    return "pytest" in joined or "py.test" in joined


def _command_looks_like_compileall(command: list[str]) -> bool:
    joined = " ".join(str(c) for c in command).lower()
    return "compileall" in joined


def _has_json_assertions(step: VerificationStep) -> bool:
    return bool(step.expect_json_keys or step.expect_json_path or step.expect_min_items)


def _get_json_path(obj: Any, path: str) -> Any:
    cur: Any = obj
    for part in path.split("."):
        if not part:
            continue
        if cur is None:
            raise KeyError(path)
        if isinstance(cur, list):
            cur = cur[int(part)]
        else:
            cur = cur[part]  # type: ignore[index]
    return cur


def _check_json_assertions(step: VerificationStep, body: Any) -> tuple[bool, str]:
    if not isinstance(body, dict):
        return False, f"expected JSON object, got {type(body).__name__}"
    for k in step.expect_json_keys:
        if k not in body:
            return False, f"missing top-level key {k!r}"
    for path, pattern in step.expect_json_path.items():
        try:
            val = _get_json_path(body, path)
        except (KeyError, IndexError, TypeError, ValueError) as e:
            return False, f"path {path!r}: {e}"
        if pattern == "*":
            continue
        if isinstance(pattern, str) and pattern.startswith("regex:"):
            rx = re.compile(pattern[6:])
            if not rx.search(str(val)):
                return False, f"path {path!r}: value {val!r} did not match {pattern!r}"
        elif str(val) != pattern:
            return False, f"path {path!r}: expected {pattern!r}, got {val!r}"
    for arr_key, min_len in step.expect_min_items.items():
        if arr_key not in body:
            return False, f"missing key {arr_key!r} for expect_min_items"
        seq = body[arr_key]
        if not isinstance(seq, list):
            return False, f"expect_min_items.{arr_key}: not a list"
        if len(seq) < min_len:
            return False, f"{arr_key} has len {len(seq)}, expected >= {min_len}"
    return True, ""


def _probe_status_only(url: str, expect_status: int) -> tuple[bool, str]:
    req = urllib.request.Request(url, method="GET")
    with urllib.request.urlopen(req, timeout=15) as resp:  # noqa: S310
        if resp.status != expect_status:
            return False, f"HTTP {resp.status} (expected {expect_status})"
    return True, ""


def _probe_json_url(url: str, step: VerificationStep) -> tuple[bool, str]:
    req = urllib.request.Request(url, method="GET")
    with urllib.request.urlopen(req, timeout=15) as resp:  # noqa: S310
        status = resp.status
        raw = resp.read().decode("utf-8", errors="replace")
    if status != step.expect_status:
        return False, f"HTTP {status} (expected {step.expect_status})"
    if not _has_json_assertions(step):
        return True, ""
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        return False, f"invalid JSON: {e}"
    ok, err = _check_json_assertions(step, data)
    if not ok:
        return False, err
    return True, ""


def _run_http_boot_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    if not step.command:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="http_boot step requires a non-empty `command`",
        )
    if not step.health_url:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="http_boot step requires a `health_url`",
        )

    effective_health_url = step.health_url
    effective_extra_health_urls = list(step.extra_health_urls)
    effective_command = list(step.command)
    effective_env = dict(env)
    port_remap_note = ""

    # Pre-flight: if something is already serving the health URL before we
    # even start the child, we would get a false positive (a zombie server
    # from an earlier run answers 200 even though our command exits). Fail
    # loudly so the agent sees the collision and can pick a free port.
    pre_existing = _probe_health_url(effective_health_url, timeout=1.0)
    if pre_existing == "alive":
        if _kill_stale_listener_for_url(effective_health_url, workspace_path=cwd):
            time.sleep(0.5)
            pre_existing = _probe_health_url(effective_health_url, timeout=1.0)
        if pre_existing != "alive":
            log.info(
                "Killed stale listener for %s before http_boot", effective_health_url
            )
        else:
            remapped = _remap_http_boot_port_for_collision(
                health_url=effective_health_url,
                extra_health_urls=effective_extra_health_urls,
                command=effective_command,
                env=effective_env,
            )
            if remapped is not None:
                (
                    effective_health_url,
                    effective_extra_health_urls,
                    effective_command,
                    effective_env,
                    port_remap_note,
                ) = remapped
                log.info(
                    "Port collision on %s; remapped http_boot to %s",
                    step.health_url,
                    effective_health_url,
                )
                if _probe_health_url(effective_health_url, timeout=1.0) == "alive":
                    return VerificationStepResult(
                        step=step,
                        status=VerificationStatus.ERROR,
                        duration_seconds=0.0,
                        summary=(
                            f"Port collision persists after remap: "
                            f"{effective_health_url} is already busy before "
                            f"{_fmt_cmd(effective_command)}."
                        ),
                        error="pre-existing listener on remapped health_url",
                    )
            else:
                return VerificationStepResult(
                    step=step,
                    status=VerificationStatus.ERROR,
                    duration_seconds=0.0,
                    summary=(
                        f"Port collision: something is already answering "
                        f"{step.health_url} before {_fmt_cmd(step.command)} was "
                        f"started. Kill the stale process or change the port."
                    ),
                    error="pre-existing listener on health_url",
                )

    start = time.time()
    stdout_buf: list[str] = []
    stderr_buf: list[str] = []

    creationflags = 0
    preexec_fn = None
    if sys.platform == "win32":
        creationflags = subprocess.CREATE_NEW_PROCESS_GROUP  # type: ignore[attr-defined]
    else:
        preexec_fn = os.setsid  # type: ignore[assignment]

    try:
        proc = subprocess.Popen(  # noqa: S603 - intentional subprocess
            effective_command,
            cwd=str(cwd),
            env=effective_env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=creationflags,
            preexec_fn=preexec_fn,
        )
    except FileNotFoundError as exc:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            duration_seconds=time.time() - start,
            error=f"executable not found: {exc}",
            summary=f"Cannot launch {_fmt_cmd(step.command)}: {exc}",
        )

    # Drain child output into background buffers so a chatty server cannot
    # fill the OS pipe and deadlock. We read the buffers after the process
    # is torn down in the finally-block.
    stdout_thread = threading.Thread(
        target=_drain_stream, args=(proc.stdout, stdout_buf), daemon=True
    )
    stderr_thread = threading.Thread(
        target=_drain_stream, args=(proc.stderr, stderr_buf), daemon=True
    )
    stdout_thread.start()
    stderr_thread.start()

    healthy = False
    last_error: str = ""
    exit_code_snapshot: int | None = None
    try:
        deadline = start + max(1, step.startup_timeout_seconds)
        while time.time() < deadline:
            if proc.poll() is not None:
                # Child exited before becoming healthy.
                break
            try:
                req = urllib.request.Request(effective_health_url, method="GET")
                with urllib.request.urlopen(req, timeout=3) as resp:  # noqa: S310
                    if 200 <= resp.status < 300:
                        # Guard against a zombie listener on the same port
                        # that started answering mid-flight after our child
                        # died (race vs. the poll() check above).
                        if proc.poll() is None:
                            healthy = True
                            break
                        last_error = (
                            "health URL answered after the child process "
                            "exited (stale listener, not our server)"
                        )
                        break
                    last_error = f"HTTP {resp.status}"
            except urllib.error.URLError as exc:
                last_error = f"URLError: {exc.reason}"
            except Exception as exc:  # noqa: BLE001
                last_error = f"{type(exc).__name__}: {exc}"
            time.sleep(0.5)

        exit_code_snapshot = proc.poll()

        # Extra HTTP assertions must run while the child is still alive.
        # (The ``finally`` block below terminates it.)
        if healthy:
            url_list: list[str] = []
            seen_u: set[str] = set()
            for u in (effective_health_url, *effective_extra_health_urls):
                u = str(u).strip()
                if u and u not in seen_u:
                    seen_u.add(u)
                    url_list.append(u)
            assert_errors: list[str] = []
            for u in url_list:
                try:
                    if _has_json_assertions(step):
                        ok, err = _probe_json_url(u, step)
                    else:
                        ok, err = _probe_status_only(u, step.expect_status)
                    if not ok:
                        assert_errors.append(f"{u}: {err}")
                except urllib.error.HTTPError as exc:
                    assert_errors.append(f"{u}: HTTPError {exc.code} {exc.reason}")
                except urllib.error.URLError as exc:
                    assert_errors.append(f"{u}: URLError {exc.reason}")
                except Exception as exc:  # noqa: BLE001
                    assert_errors.append(f"{u}: {type(exc).__name__}: {exc}")
            duration = time.time() - start
            stdout_text = "".join(stdout_buf)
            stderr_text = "".join(stderr_buf)
            if assert_errors:
                return VerificationStepResult(
                    step=step,
                    status=VerificationStatus.FAILED,
                    duration_seconds=duration,
                    summary=(
                        f"HTTP assertions failed after boot "
                        f"({_fmt_cmd(step.command)}): " + "; ".join(assert_errors[:5])
                    ),
                    stdout=stdout_text,
                    stderr=stderr_text,
                    error="; ".join(assert_errors)[:2000],
                )
            return VerificationStepResult(
                step=step,
                status=VerificationStatus.PASSED,
                duration_seconds=duration,
                summary=(
                    f"Service booted: {_fmt_cmd(effective_command)} responded "
                    f"200 at {effective_health_url} in {duration:.1f}s"
                    + (f" ({port_remap_note})" if port_remap_note else "")
                ),
                stdout=stdout_text,
                stderr=stderr_text,
            )
    finally:
        _terminate_process(proc)
        # Make sure the drain threads finish reading whatever the process
        # wrote on its way out so the buffers are complete.
        stdout_thread.join(timeout=5)
        stderr_thread.join(timeout=5)

    duration = time.time() - start
    stdout_text = "".join(stdout_buf)
    stderr_text = "".join(stderr_buf)

    return VerificationStepResult(
        step=step,
        status=VerificationStatus.FAILED,
        exit_code=exit_code_snapshot,
        duration_seconds=duration,
        summary=(
            f"Service did not become healthy at {effective_health_url} "
            f"within {step.startup_timeout_seconds}s"
            + (f" ({port_remap_note})" if port_remap_note else "")
        ),
        stdout=stdout_text,
        stderr=stderr_text,
        error=last_error,
    )


def _remap_http_boot_port_for_collision(
    *,
    health_url: str,
    extra_health_urls: list[str],
    command: list[str],
    env: dict[str, str],
) -> tuple[str, list[str], list[str], dict[str, str], str] | None:
    """Try to remap a colliding localhost health URL to a free port.

    Returns remapped ``(health_url, extra_urls, command, env, note)`` or
    ``None`` when remap is not possible.
    """
    parsed = urllib.parse.urlparse(health_url)
    if parsed.hostname not in {"127.0.0.1", "localhost", "::1"}:
        return None
    if not parsed.port:
        return None
    free_port = _reserve_free_local_port()
    if free_port is None:
        return None

    def _replace_port(url: str) -> str:
        p = urllib.parse.urlparse(url)
        if p.hostname not in {"127.0.0.1", "localhost", "::1"} or not p.port:
            return url
        host = p.hostname
        if ":" in host and not host.startswith("["):
            host = f"[{host}]"
        netloc = f"{host}:{free_port}"
        return urllib.parse.urlunparse(
            (p.scheme, netloc, p.path, p.params, p.query, p.fragment)
        )

    remapped_health = _replace_port(health_url)
    remapped_extra = [_replace_port(u) for u in extra_health_urls]
    remapped_env = dict(env)
    remapped_env["PORT"] = str(free_port)
    remapped_env["UVICORN_PORT"] = str(free_port)
    remapped_command = _inject_or_replace_port_arg(command, free_port)
    note = f"auto-remapped port to {free_port}"
    return remapped_health, remapped_extra, remapped_command, remapped_env, note


def _reserve_free_local_port() -> int | None:
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            return int(s.getsockname()[1])
    except OSError:
        return None


def _inject_or_replace_port_arg(command: list[str], port: int) -> list[str]:
    cmd = list(command)
    for i, token in enumerate(cmd):
        low = str(token).lower()
        if low == "--port" and i + 1 < len(cmd):
            cmd[i + 1] = str(port)
            return cmd
        if low.startswith("--port="):
            cmd[i] = f"--port={port}"
            return cmd
    cmd_low = [str(t).lower() for t in cmd]
    if "uvicorn" in cmd_low or any("uvicorn" in t for t in cmd_low):
        cmd.extend(["--port", str(port)])
    return cmd


def _post_json(
    url: str, payload: dict[str, Any], *, timeout: float = 20.0
) -> tuple[int, str]:
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=data,
        method="POST",
        headers={"Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:  # noqa: S310
            return resp.status, resp.read().decode("utf-8", errors="replace")
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace")
        return exc.code, body or f"HTTPError {exc.code}: {exc.reason}"


def _run_behavioral_http_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    if not step.command:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="behavioral_http requires command",
        )
    request_url = step.request_url or (step.health_url.rstrip("/") + "/generate")
    payloads = step.request_payloads or [{"input": "alpha"}, {"input": "beta"}]
    if len(payloads) < 2:
        payloads = (
            [payloads[0], {"input": "beta"}]
            if payloads
            else [{"input": "alpha"}, {"input": "beta"}]
        )

    start = time.time()
    stdout_buf: list[str] = []
    stderr_buf: list[str] = []
    creationflags = (
        subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == "win32" else 0
    )  # type: ignore[attr-defined]
    preexec_fn = None if sys.platform == "win32" else os.setsid  # type: ignore[assignment]
    try:
        proc = subprocess.Popen(
            list(step.command),
            cwd=str(cwd),
            env=env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=creationflags,
            preexec_fn=preexec_fn,
        )
    except FileNotFoundError as exc:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error=f"executable not found: {exc}",
        )

    stdout_thread = threading.Thread(
        target=_drain_stream, args=(proc.stdout, stdout_buf), daemon=True
    )
    stderr_thread = threading.Thread(
        target=_drain_stream, args=(proc.stderr, stderr_buf), daemon=True
    )
    stdout_thread.start()
    stderr_thread.start()
    try:
        deadline = start + max(1, step.startup_timeout_seconds)
        last_error = ""
        while time.time() < deadline:
            if proc.poll() is not None:
                last_error = f"process exited early rc={proc.returncode}"
                break
            try:
                req = urllib.request.Request(step.health_url, method="GET")
                with urllib.request.urlopen(req, timeout=3) as resp:  # noqa: S310
                    if 200 <= resp.status < 300:
                        break
            except Exception as exc:  # noqa: BLE001
                last_error = f"{type(exc).__name__}: {exc}"
            time.sleep(0.5)
        else:
            return VerificationStepResult(
                step=step,
                status=VerificationStatus.FAILED,
                duration_seconds=time.time() - start,
                summary=f"Service did not become healthy before behavioral probe: {last_error}",
                stdout="".join(stdout_buf),
                stderr="".join(stderr_buf),
                error=last_error,
            )

        bodies: list[str] = []
        statuses: list[int] = []
        mock_hits: list[str] = []
        for payload in payloads[:2]:
            status, body = _post_json(request_url, payload)
            statuses.append(status)
            body = body.strip()
            bodies.append(body)
            mock_hits.extend(_mock_scaffold_hits(body))
        semantic_ok, semantic_reason = _semantic_http_outputs_match_inputs(
            payloads[:2], bodies
        )
        ok = (
            all(s == step.expect_status for s in statuses)
            and all(bodies)
            and bodies[0] != bodies[1]
            and not mock_hits
            and semantic_ok
        )
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.PASSED if ok else VerificationStatus.FAILED,
            duration_seconds=time.time() - start,
            summary=(
                f"POST {request_url} statuses={statuses}; "
                f"non_empty={all(bool(b) for b in bodies)}; "
                f"different_outputs={len(bodies) >= 2 and bodies[0] != bodies[1]}; "
                f"semantic_input_match={semantic_ok}"
            ),
            stdout="\n--- request A ---\n"
            + json.dumps(payloads[0], ensure_ascii=False, indent=2)
            + "\n--- response A ---\n"
            + (bodies[0] if bodies else "")
            + "\n--- request B ---\n"
            + json.dumps(payloads[1], ensure_ascii=False, indent=2)
            + "\n--- response B ---\n"
            + (bodies[1] if len(bodies) > 1 else ""),
            stderr="".join(stderr_buf),
            error=""
            if ok
            else (
                "mock/placeholder output detected: " + ", ".join(sorted(set(mock_hits)))
                if mock_hits
                else semantic_reason
                if not semantic_ok
                else "behavioral outputs were empty, same, or non-200"
            ),
        )
    except Exception as exc:  # noqa: BLE001
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.FAILED,
            duration_seconds=time.time() - start,
            stdout="".join(stdout_buf),
            stderr="".join(stderr_buf),
            error=f"{type(exc).__name__}: {exc}",
        )
    finally:
        _terminate_process(proc)
        stdout_thread.join(timeout=5)
        stderr_thread.join(timeout=5)


def _run_input_sensitivity_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    if not step.command:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="input_sensitivity_check requires command",
        )
    a = _run_shell_step(
        step, cwd, {**env, "UMBRELLA_VERIFY_INPUT": "alpha verification input"}
    )
    b = _run_shell_step(
        step, cwd, {**env, "UMBRELLA_VERIFY_INPUT": "beta verification input"}
    )
    combined = (a.stdout.strip(), b.stdout.strip())
    passed = (
        a.status == VerificationStatus.PASSED
        and b.status == VerificationStatus.PASSED
        and combined[0]
        and combined[0] != combined[1]
    )
    return VerificationStepResult(
        step=step,
        status=VerificationStatus.PASSED if passed else VerificationStatus.FAILED,
        summary=f"input sensitivity: run_a={a.status.value}, run_b={b.status.value}, different_outputs={combined[0] != combined[1]}",
        stdout=f"--- A ---\n{combined[0]}\n--- B ---\n{combined[1]}",
        stderr=(a.stderr + "\n" + b.stderr).strip(),
        error="" if passed else "outputs were empty, identical, or one command failed",
    )


def _run_pptx_diff_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    del env
    pattern = step.output_glob or "*.pptx"
    pptx_files = [
        p for p in cwd.glob(pattern) if not p.name.lower().startswith("template")
    ]
    template_files = [p for p in cwd.glob("template*.pptx")]
    if not pptx_files:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.FAILED,
            summary=f"No generated pptx matching {pattern}",
        )
    if not template_files:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.PASSED,
            summary=f"Found generated pptx: {pptx_files[0].name}",
        )
    generated_path = pptx_files[0]
    template_path = template_files[0]
    generated = generated_path.read_bytes()
    template = template_path.read_bytes()
    text_result = _pptx_text_diff_summary(generated_path, template_path)
    passed = generated != template and len(generated) > 0 and text_result[0]
    return VerificationStepResult(
        step=step,
        status=VerificationStatus.PASSED if passed else VerificationStatus.FAILED,
        summary=(
            f"Generated {generated_path.name} differs from template {template_path.name}: "
            f"bytes={generated != template}; {text_result[1]}"
        ),
        error=""
        if passed
        else "generated pptx is byte-identical, empty, or has no meaningful text changes",
    )


def _mock_scaffold_hits(text: str) -> list[str]:
    return mock_scaffold_hits(text)


def _semantic_http_outputs_match_inputs(
    payloads: list[dict[str, Any]],
    bodies: list[str],
) -> tuple[bool, str]:
    """Require each output to retain topic words from its own input."""
    if len(payloads) < 2 or len(bodies) < 2:
        return False, "semantic_input_match requires at least two payloads and bodies"
    body_a = bodies[0].lower()
    body_b = bodies[1].lower()
    tokens_a = _payload_topic_tokens(payloads[0])
    tokens_b = _payload_topic_tokens(payloads[1])
    if not tokens_a or not tokens_b:
        return True, ""
    hits_a = [token for token in tokens_a if token in body_a]
    hits_b = [token for token in tokens_b if token in body_b]
    cross_a = [token for token in tokens_a if token in body_b]
    cross_b = [token for token in tokens_b if token in body_a]
    own_a = bool(hits_a)
    own_b = bool(hits_b)
    cross_dominates = len(cross_a) > len(hits_a) or len(cross_b) > len(hits_b)
    if own_a and own_b and not cross_dominates:
        return True, ""
    return (
        False,
        "behavioral outputs do not preserve topic tokens from their corresponding inputs",
    )


def _payload_topic_tokens(payload: Any) -> list[str]:
    text = json.dumps(payload, ensure_ascii=False, default=str)
    stop = {
        "input",
        "text",
        "content",
        "data",
        "the",
        "and",
        "with",
        "from",
        "this",
        "that",
        "для",
        "или",
        "как",
        "это",
        "что",
        "новости",
    }
    tokens: list[str] = []
    for token in re.findall(r"[A-Za-zА-Яа-я0-9]{4,}", text.lower()):
        if token not in stop and token not in tokens:
            tokens.append(token)
    return tokens[:12]


def _pptx_text_diff_summary(
    generated_path: Path, template_path: Path
) -> tuple[bool, str]:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        return True, f"text diff skipped: python-pptx unavailable ({exc})"

    def _texts(path: Path) -> list[str]:
        prs = Presentation(str(path))
        values: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    values.append(text.strip())
        return values

    try:
        generated = _texts(generated_path)
        template = _texts(template_path)
    except Exception as exc:  # noqa: BLE001
        return False, f"text diff failed: {type(exc).__name__}: {exc}"
    changed = sum(1 for left, right in zip(generated, template) if left != right)
    extra = abs(len(generated) - len(template))
    meaningful = bool(generated) and (changed + extra) > 0
    return (
        meaningful,
        f"text_changed={changed}; text_count={len(generated)}; template_text_count={len(template)}",
    )


def _kill_stale_listener_for_url(url: str, *, workspace_path: Path) -> bool:
    """Best-effort kill for a stale local listener owned by this workspace."""
    try:
        parsed = urllib.parse.urlparse(url)
        host = parsed.hostname
        port = parsed.port
        if not port or host not in {"127.0.0.1", "localhost", "::1"}:
            return False
    except Exception:
        return False
    try:
        import psutil  # type: ignore[import-not-found]
    except Exception:
        return False
    workspace_str = str(workspace_path.resolve()).lower()
    killed = False
    try:
        for conn in psutil.net_connections(kind="inet"):
            if not conn.laddr or conn.laddr.port != port or not conn.pid:
                continue
            try:
                proc = psutil.Process(conn.pid)
                cmdline = " ".join(proc.cmdline()).lower()
                cwd = (proc.cwd() or "").lower()
                if workspace_str not in cmdline and workspace_str not in cwd:
                    continue
                proc.kill()
                killed = True
            except Exception:
                log.debug(
                    "Failed to kill stale listener pid=%s", conn.pid, exc_info=True
                )
    except Exception:
        log.debug("Failed to inspect local listeners for %s", url, exc_info=True)
    return killed


def _probe_health_url(url: str, *, timeout: float = 1.0) -> str:
    """Fire a single HEAD/GET at ``url`` and classify the response.

    Returns:
        - ``"alive"``  — something responded with 2xx/3xx/4xx/5xx.
        - ``"dead"``   — connection refused / DNS failure / IO error.
        - ``"error"``  — malformed URL or unexpected exception.
    """
    try:
        req = urllib.request.Request(url, method="GET")
        with urllib.request.urlopen(req, timeout=timeout) as resp:  # noqa: S310
            _ = resp.status  # any HTTP response counts as alive
            return "alive"
    except urllib.error.HTTPError:
        # A 4xx/5xx still means *something* is listening.
        return "alive"
    except urllib.error.URLError:
        return "dead"
    except Exception:  # noqa: BLE001
        return "error"


def _drain_stream(stream: IO[str] | None, buf: list[str]) -> None:
    """Read all output from ``stream`` into ``buf`` until EOF.

    Used as a target for background threads that keep the child's pipes
    drained (a blocked pipe deadlocks the child).
    """
    if stream is None:
        return
    try:
        for chunk in iter(lambda: stream.read(4096), ""):
            if not chunk:
                break
            buf.append(chunk)
    except Exception:  # noqa: BLE001
        log.debug("Failed to drain child stream", exc_info=True)
    finally:
        try:
            stream.close()
        except Exception:  # noqa: BLE001
            pass


def _terminate_process(proc: subprocess.Popen) -> None:
    if proc.poll() is not None:
        return
    try:
        if sys.platform == "win32":
            proc.send_signal(signal.CTRL_BREAK_EVENT)  # type: ignore[attr-defined]
        else:
            os.killpg(os.getpgid(proc.pid), signal.SIGTERM)
    except Exception:  # noqa: BLE001
        proc.terminate()

    try:
        proc.wait(timeout=5)
    except subprocess.TimeoutExpired:
        try:
            if sys.platform != "win32":
                os.killpg(os.getpgid(proc.pid), signal.SIGKILL)
        except Exception:  # noqa: BLE001
            pass
        proc.kill()


# ---------------------------------------------------------------------------
# Import-check step
# ---------------------------------------------------------------------------


def _run_import_step(
    step: VerificationStep, cwd: Path, env: dict[str, str]
) -> VerificationStepResult:
    if step.command:
        cmd = list(step.command)
    elif step.module:
        cmd = [sys.executable, "-c", f"import {step.module}"]
    else:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.ERROR,
            error="import_check step requires either `module` or a `command`",
        )

    start = time.time()
    try:
        proc = subprocess.run(  # noqa: S603
            cmd,
            cwd=str(cwd),
            env=env,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=max(1, step.timeout_seconds),
            check=False,
        )
    except subprocess.TimeoutExpired:
        return VerificationStepResult(
            step=step,
            status=VerificationStatus.FAILED,
            duration_seconds=time.time() - start,
            summary=f"import_check timed out: {_fmt_cmd(cmd)}",
            error="timeout",
        )

    duration = time.time() - start
    status = (
        VerificationStatus.PASSED if proc.returncode == 0 else VerificationStatus.FAILED
    )
    return VerificationStepResult(
        step=step,
        status=status,
        exit_code=proc.returncode,
        duration_seconds=duration,
        stdout=proc.stdout or "",
        stderr=proc.stderr or "",
        summary=(
            f"import check `{step.module or _fmt_cmd(cmd)}` -> exit {proc.returncode}"
        ),
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _fmt_cmd(cmd: Iterable[str]) -> str:
    try:
        return " ".join(shlex.quote(str(c)) for c in cmd)
    except Exception:  # noqa: BLE001
        return " ".join(str(c) for c in cmd)
